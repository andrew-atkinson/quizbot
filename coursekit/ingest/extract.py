"""Source document -> raw text. Deterministic, offline, no model.

Each supported type has a small extractor. PDF/PPTX use lazily-imported parsers (`pypdf`,
`python-pptx`); Word (`.docx`) and OpenDocument (`.odt`) are zip-of-XML, so they are read with the
**stdlib only** (`zipfile` + `xml.etree`) — no extra dependency. Extraction is deliberately shallow:
it pulls the text out and does light cosmetic cleanup; turning messy extraction into a teaching-ready
document is the optional shaping pass (`ingest.shape`), not this layer's job.

Known limitations (text only — see agent/todo.md for the future-iteration list):
  - **Text, not visuals.** Images, diagrams, charts, and photos in slides/PDFs are ignored; nothing
    describes them. (The video transcriber has a vision stage; ingest could grow one — the heavy axis.)
  - **No OCR.** A scanned PDF (an image of a page) yields no text.
  - **Shallow layout.** Multi-column PDFs, tables, and complex layouts extract as linear text and can
    interleave; the shaping pass mitigates but does not reconstruct structure.
  - **PPTX** now includes slide text + speaker notes, but not chart data, SmartArt, or embedded objects.
  - **`.doc`** (legacy binary Word) and **direct Google Slides/Docs** import are unsupported — convert
    or export first.
"""

import re
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

_MULTISPACE = re.compile(r"[ \t]+")
_MANY_BLANKS = re.compile(r"\n{3,}")


def _clean(text: str) -> str:
    """Cosmetic only: trim trailing spaces per line, collapse runs of blank lines. Preserve the
    words — the shaping pass (or the generator) does the real structuring."""
    lines = [_MULTISPACE.sub(" ", ln).rstrip() for ln in (text or "").splitlines()]
    return _MANY_BLANKS.sub("\n\n", "\n".join(lines)).strip()


def _text(path: Path) -> str:
    return _clean(path.read_text(encoding="utf-8", errors="replace"))


def _pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return _clean("\n\n".join((page.extract_text() or "") for page in reader.pages))


def _pptx(path: Path) -> str:
    """Slide text plus **speaker notes** — on lecture decks the substance is usually in the notes,
    with the slide itself just a headline, so both are pulled (slide text first, then its notes)."""
    from pptx import Presentation
    blocks = []
    for slide in Presentation(str(path)).slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if t:
                    texts.append(t)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                texts.append(notes)
        if texts:
            blocks.append("\n".join(texts))
    return _clean("\n\n".join(blocks))


_DOCX_T = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"
_DOCX_P = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p"
_ODT_NS = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"


def _docx(path: Path) -> str:
    """Word .docx = a zip; the body text lives in word/document.xml as <w:t> runs inside <w:p>
    paragraphs. Stdlib only."""
    with zipfile.ZipFile(path) as z:
        root = ET.fromstring(z.read("word/document.xml"))
    paras = ["".join(t.text or "" for t in p.iter(_DOCX_T)) for p in root.iter(_DOCX_P)]
    return _clean("\n".join(paras))


def _odt(path: Path) -> str:
    """OpenDocument .odt = a zip; text is in content.xml as <text:p>/<text:h> (with nested
    <text:span>). itertext() gathers the runs; iter() yields document order. Stdlib only."""
    with zipfile.ZipFile(path) as z:
        root = ET.fromstring(z.read("content.xml"))
    blocks = [f"{_ODT_NS}p", f"{_ODT_NS}h"]
    paras = ["".join(el.itertext()) for el in root.iter() if el.tag in blocks]
    return _clean("\n".join(paras))


# suffix -> extractor. The keys are also the "supported source" set.
_EXTRACTORS = {
    ".pdf": _pdf,
    ".pptx": _pptx,
    ".docx": _docx,
    ".odt": _odt,
    ".txt": _text,
    ".md": _text,
    ".markdown": _text,
}

SUPPORTED_SUFFIXES = frozenset(_EXTRACTORS)


def is_supported(path) -> bool:
    return Path(path).suffix.lower() in _EXTRACTORS


def extract_text(path) -> str:
    """The document's text, cleaned. Raises ValueError for an unsupported type."""
    path = Path(path)
    fn = _EXTRACTORS.get(path.suffix.lower())
    if fn is None:
        raise ValueError(
            f"unsupported source type '{path.suffix}' ({path.name}); "
            f"supported: {', '.join(sorted(SUPPORTED_SUFFIXES))}")
    return fn(path)
