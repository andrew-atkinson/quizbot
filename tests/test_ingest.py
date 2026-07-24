"""Document ingest: source -> output/week-N.md. Offline — extraction is deterministic, and the
shaping pass is exercised through a fake Provider (no network, no real model)."""

import inspect
import io
import zipfile
from pathlib import Path

import pytest

from coursekit.ingest import extract, ingest


# ------------------------------------------------------------- extraction

def _make_pdf(text: str) -> bytes:
    """A minimal single-page PDF with one text line and a correct xref table — enough for pypdf to
    extract `text`. (No PDF-writer lib is available, so we assemble the bytes.)"""
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
    ]
    content = b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
    objs.append(b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content
                + b"\nendstream")
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, 1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + body + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (b"trailer\n<< /Size " + str(len(objs) + 1).encode() + b" /Root 1 0 R >>\n"
            b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF")
    return bytes(out)


def test_extract_pdf_real(tmp_path):
    f = tmp_path / "week-4.pdf"
    f.write_bytes(_make_pdf("Hello Week Four Randomness"))
    assert "Randomness" in extract.extract_text(f)


def test_extract_pptx_roundtrip(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # "Title Only"
    slide.shapes.title.text = "Exposure Triangle"
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(2))
    tb.text_frame.text = "Aperture, shutter, ISO"
    f = tmp_path / "week-2.pptx"
    prs.save(str(f))

    text = extract.extract_text(f)
    assert "Exposure Triangle" in text and "Aperture, shutter, ISO" in text


def test_extract_pptx_includes_speaker_notes(tmp_path):
    # lecture decks keep the substance in the notes; the slide is just a headline
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Aperture"
    slide.notes_slide.notes_text_frame.text = "Aperture controls depth of field and light."
    f = tmp_path / "week-2.pptx"
    prs.save(str(f))

    text = extract.extract_text(f)
    assert "Aperture" in text
    assert "controls depth of field" in text     # the speaker-notes content came through


def _make_docx(paras: list[str]) -> bytes:
    """A minimal .docx: a zip whose word/document.xml holds <w:p> paragraphs of <w:t> runs."""
    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    body = "".join(f"<w:p><w:r><w:t>{p}</w:t></w:r></w:p>" for p in paras)
    doc = f'<?xml version="1.0"?><w:document xmlns:w="{W}"><w:body>{body}</w:body></w:document>'
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("word/document.xml", doc)
    return buf.getvalue()


def _make_odt(paras: list[str]) -> bytes:
    """A minimal .odt: a zip whose content.xml holds <text:p> paragraphs (one with a nested span)."""
    T = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    O = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
    body = "".join(f"<text:p>{p}</text:p>" for p in paras)
    content = (f'<?xml version="1.0"?><office:document-content xmlns:office="{O}" '
               f'xmlns:text="{T}"><office:body><office:text>{body}</office:text>'
               f'</office:body></office:document-content>')
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("content.xml", content)
    return buf.getvalue()


def test_extract_docx(tmp_path):
    f = tmp_path / "week-6.docx"
    f.write_bytes(_make_docx(["Depth of field explained", "Aperture and focal length"]))
    text = extract.extract_text(f)
    assert "Depth of field explained" in text and "Aperture and focal length" in text


def test_extract_odt_gathers_nested_spans(tmp_path):
    f = tmp_path / "week-7.odt"
    f.write_bytes(_make_odt(["Colour theory intro", "Warm <text:span>and cool</text:span> tones"]))
    text = extract.extract_text(f)
    assert "Colour theory intro" in text
    assert "Warm and cool tones" in text        # nested span text is gathered, not dropped


def test_extract_txt_cleans_whitespace(tmp_path):
    t = tmp_path / "a.txt"
    t.write_text("line one\n\n\n\nline two   \n", encoding="utf-8")
    assert extract.extract_text(t) == "line one\n\nline two"   # 3+ blanks collapsed, lines rstripped


def test_extract_md_passthrough(tmp_path):
    m = tmp_path / "b.md"
    m.write_text("# Heading\nbody text", encoding="utf-8")
    assert "# Heading" in extract.extract_text(m)


def test_unsupported_type_raises():
    with pytest.raises(ValueError, match="unsupported"):
        extract.extract_text(Path("notes.rtf"))       # .rtf is not in the supported set


# ------------------------------------------------------------- week mapping

def test_plan_weeks_uses_filename_week_numbers(tmp_path):
    plan = ingest.plan_weeks([tmp_path / "week-3.pdf", tmp_path / "week-1.pdf"])
    assert [s for s, _ in plan] == ["week-1", "week-3"]      # sorted, keyed by name


def test_plan_weeks_enumerates_when_unnumbered(tmp_path):
    plan = ingest.plan_weeks([tmp_path / "intro.pdf", tmp_path / "basics.pdf"])
    assert [s for s, _ in plan] == ["week-1", "week-2"]      # sorted name order
    assert plan[0][1].name == "basics.pdf"


def test_plan_weeks_rejects_a_duplicate_week(tmp_path):
    with pytest.raises(ValueError, match="week-3"):
        ingest.plan_weeks([tmp_path / "week-3.pdf", tmp_path / "week-3.txt"])


# ------------------------------------------------------------- orchestration

class _FakeProvider:
    def __init__(self, reply="# shaped\nclean doc"):
        self.calls, self._reply = [], reply

    def chat(self, *, model, messages, temperature=None, max_tokens=None):
        self.calls.append({"model": model, "messages": messages})
        return self._reply


def test_ingest_raw_writes_week_md_without_a_model(tmp_path):
    (tmp_path / "week-2.txt").write_text("some reading text", encoding="utf-8")
    out = ingest.ingest(tmp_path, raw=True)                  # no provider at all
    assert len(out) == 1
    src, dest = out[0]
    assert dest == tmp_path / "output" / "week-2.md"
    assert "some reading text" in dest.read_text(encoding="utf-8")


def test_ingest_shapes_each_doc_through_the_provider(tmp_path):
    (tmp_path / "week-3.txt").write_text("raw   messy    text", encoding="utf-8")
    prov = _FakeProvider("# Week 3\nclean")
    out = ingest.ingest(tmp_path, raw=False, provider=prov, model="m")

    dest = out[0][1]
    assert dest.name == "week-3.md" and dest.read_text(encoding="utf-8").strip() == "# Week 3\nclean"
    msgs = prov.calls[0]["messages"]
    assert msgs[0]["role"] == "system" and "FAITHFUL cleanup" in msgs[0]["content"]
    assert "raw messy text" in msgs[1]["content"]           # the cleaned extraction was sent


def test_shaping_without_a_provider_errors(tmp_path):
    (tmp_path / "week-1.txt").write_text("x", encoding="utf-8")
    with pytest.raises(ValueError, match="provider"):
        ingest.ingest(tmp_path, raw=False, provider=None)


def test_ingest_no_supported_docs_is_empty(tmp_path):
    (tmp_path / "notes.rtf").write_text("x", encoding="utf-8")   # unsupported type
    assert ingest.ingest(tmp_path, raw=True) == []


def test_ingest_writes_into_the_vtconfig_course_output(tmp_path):
    # a real course root: the week docs land in <root>/output/, where discover.find_units looks
    root = tmp_path / "course"
    (root / ".vtconfig").mkdir(parents=True)
    (root / "readings").mkdir()
    (root / "readings" / "week-5.txt").write_text("exposure and aperture", encoding="utf-8")
    out = ingest.ingest(root / "readings", raw=True)
    assert out[0][1] == root / "output" / "week-5.md"


# ------------------------------------------------------------- offline guarantee

def test_ingest_makes_no_network_calls():
    # same spirit as the no-URL guardrail: the ingest path is local-only, never fetches. Check for
    # network *libraries* (not literal "http://", which appears legitimately in XML namespace URIs).
    src = inspect.getsource(ingest) + inspect.getsource(extract)
    for bad in ("requests", "urllib", "httpx", "socket", "urlopen"):
        assert bad not in src, f"ingest must stay offline, found {bad!r}"
